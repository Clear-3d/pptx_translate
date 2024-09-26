from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx import Presentation

# 创建Presentation对象
prs = Presentation()

# 添加新的幻灯片
slide_layout = prs.slide_layouts[0]  # 选择第一个幻灯片布局
slide = prs.slides.add_slide(slide_layout)


# 添加形状到幻灯片
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1))
shape.text = "Animated Shape"
shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
shape.text_frame.paragraphs[0].font.size = Pt(24)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(255, 0, 0)  # 设置形状填充颜色为红色

# 添加淡入动画效果
effect = slide.shapes._spTree.insert(2, OxmlElement('p:fade'))
effect.set('start', 'onLoad')
effect.set('dur', '1000ms')
effect.set('fadeColor', 'black')
