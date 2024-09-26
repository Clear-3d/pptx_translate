from pptx import Presentation
from pptx.util import Pt



def get_textbox_content(file_path):
    prs = Presentation(file_path)
    text_list = []
    
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        text_list.append(slide_text)
    
    print(text_list)
    return text_list

def modify_last_textbox_font(file_path):
    prs = Presentation(file_path)
    for slide in prs.slides:
        text_boxes = [shape for shape in slide.shapes if hasattr(shape, "text") and shape.text_frame]
        if text_boxes:
            last_text_box = text_boxes[-1]  # 获取最后一个文本框
            last_text_box.text_frame.paragraphs[0].font.size = Pt(2)

    prs.save('modified_presentation.pptx')

# 使用示例
file_path = 'test.pptx'
modify_last_textbox_font(file_path)
